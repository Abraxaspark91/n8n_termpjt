# -*- coding: utf-8 -*-
"""
apply_deliverable_patch.py — W3 산출물 동기화 엔진 v4 (GitHub Actions에서 실행)

v4 핵심: 시나리오 루트(scenarios/<슬러그>/deliverables) 지원 + 전역 패치 큐(patches/) + 파일명 비의존.
Claude Code(Sonnet)가 어떤 파일명·구조로 산출물을 만들어도 대응한다:
  1) <dir>/manifest.json 이 있으면 그 경로를 신뢰 (계약)
  2) 없으면 <dir> 안의 *.docx / *.xlsx / *.pptx 를 자동 탐색 (정렬 후 첫 파일)
  3) 그래도 없으면 기본 파일명으로 부트스트랩 생성 + manifest 기록 (멱등)

패치 스키마:
  deliverable_patch/v2 : {deliverable_dir?, ops[]}  ← n8n Parse Patch Spec 산출
  deliverable_patch/v1 : 구버전 루트 파일 append — 하위호환 유지

ops 어휘 (n8n 검증기와 1:1):
  docx.append_changelog / docx.append_after_heading / docx.replace_paragraph
  xlsx.append_row(RTM) / xlsx.update_by_req_id / xlsx.update_cell
  pptx.append_changelog_slide / pptx.replace_text
RTM 시트는 어떤 워크북에든 없으면 생성한다(기존 시트 불변).

상태 추출본(에이전트가 읽는 원천, 파일명 고정): <dir>/state/
  docx.md · rtm.csv · xlsx_overview.md · pptx.md
CLI: --export-state-all  → 패치 없이 deliverables/*/ 전 디렉토리 상태본만 재생성
"""
import csv
import io
import json
import shutil
import sys
from pathlib import Path

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt as PPt

ROOT = Path(__file__).resolve().parents[1]
DELIV = ROOT / "deliverables"
PATCHES = ROOT / "patches"
PENDING = PATCHES / "pending"
APPLIED = PATCHES / "applied"

RTM_HEADER = ["req_id", "일시", "요청자", "소스", "결정", "요약",
              "영향 산출물", "변경유형", "DA점수", "위키 커밋"]
MAX_OPS = 12
DEFAULTS = {"docx": "산출물_설계서.docx", "xlsx": "요구사항_추적표.xlsx", "pptx": "보고_장표.pptx"}


# ── 파일 결정: manifest → 자동 탐색 → 부트스트랩 ─────────────────────
def discover(dirpath: Path, ext: str):
    hits = sorted(p for p in dirpath.glob(f"*.{ext}") if not p.name.startswith("~$"))
    return hits[0] if hits else None


def resolve_files(dirpath: Path):
    dirpath.mkdir(parents=True, exist_ok=True)
    mf_path = dirpath / "manifest.json"
    mf = {}
    if mf_path.exists():
        try:
            mf = json.loads(mf_path.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"warn: manifest unreadable in {dirpath}: {e}")
    files, changed = {}, False
    for kind in ("docx", "xlsx", "pptx"):
        cand = mf.get(kind)
        p = (dirpath / cand) if cand and (dirpath / cand).exists() else None
        if p is None:
            p = discover(dirpath, kind)               # Sonnet이 만든 임의 파일명 대응
        if p is None:
            p = dirpath / DEFAULTS[kind]              # 부트스트랩 예정
        if mf.get(kind) != p.name:
            mf[kind] = p.name
            changed = True
        files[kind] = p
    mf.setdefault("schema", "deliverable_manifest/v1")
    mf.setdefault("rtm_sheet", "RTM")
    if changed or not mf_path.exists():
        mf_path.write_text(json.dumps(mf, ensure_ascii=False, indent=2), encoding="utf-8")
    return files, mf


def open_docx(p: Path) -> Document:
    if p.exists():
        return Document(str(p))
    doc = Document()
    doc.add_heading(p.stem, level=0)
    doc.add_paragraph("W3 파이프라인 관리 문서. 승인 요건이 자동 반영됩니다.")
    doc.add_heading("변경 이력", level=1)
    return doc


def open_xlsx(p: Path):
    wb = load_workbook(str(p)) if p.exists() else Workbook()
    if not p.exists():
        wb.active.title = "RTM"
    if "RTM" not in wb.sheetnames:                    # Sonnet 시트 구성 무관 — RTM만 보장
        wb.create_sheet("RTM")
    ws = wb["RTM"]
    if ws.max_row == 1 and all(c.value is None for c in ws[1]):
        for col, h in enumerate(RTM_HEADER, start=1):
            ws.cell(row=1, column=col, value=h)
    return wb


def open_pptx(p: Path) -> Presentation:
    if p.exists():
        return Presentation(str(p))
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = p.stem
    return prs


# ── ops 구현 (v3.2와 동일 어휘) ──────────────────────────────────────
def docx_append_changelog(doc, p, op):
    doc.add_heading(op.get("title") or f"[{p.get('req_id', '-')}] {p.get('summary', '')[:60]}", level=2)
    meta = doc.add_paragraph()
    meta.add_run(f"일시 {p.get('ts', '-')} · 요청자 {p.get('requestor', '-')} · 소스 {p.get('source', '-')} "
                 f"· 결정 {p.get('decision', '-')} · DA {p.get('da_score', '-')}").font.size = Pt(9)
    for b in op.get("bullets") or []:
        doc.add_paragraph(str(b), style="List Bullet")
    if op.get("risks"):
        doc.add_paragraph("리스크:", style="Intense Quote")
        for r in op["risks"]:
            doc.add_paragraph(str(r), style="List Bullet")
    if op.get("open_questions"):
        doc.add_paragraph("미해결 질문:", style="Intense Quote")
        for q in op["open_questions"]:
            doc.add_paragraph(str(q), style="List Bullet")
    if op.get("wiki_commit_url"):
        doc.add_paragraph(f"위키 커밋: {op['wiki_commit_url']}")
    return "applied"


def docx_append_after_heading(doc, p, op):
    target = str(op.get("heading", "")).strip()
    for para in doc.paragraphs:
        if para.style.name.lower().startswith("heading") and para.text.strip() == target:
            style = op.get("style") if op.get("style") in ("List Bullet", "Intense Quote") else None
            new = doc.add_paragraph(str(op.get("text", "")), style=style)
            para._p.addnext(new._p)
            return "applied"
    return f"skipped: heading not found ({target})"


def docx_replace_paragraph(doc, p, op):
    match = str(op.get("match", ""))
    for para in doc.paragraphs:
        if para.text == match:
            for r in list(para.runs):
                r.text = ""
            (para.runs[0] if para.runs else para.add_run()).text = str(op.get("new_text", ""))
            return "applied"
    return "skipped: paragraph not found"


def xlsx_append_row(wb, p, op):
    sheet = op.get("sheet") or "RTM"
    ws = wb[sheet] if sheet in wb.sheetnames else wb["RTM"]
    ws.append([("" if v is None else v) for v in (op.get("values") or [])])
    return "applied"


def xlsx_update_by_req_id(wb, p, op):
    ws = wb["RTM"]
    headers = [str(c.value or "") for c in ws[1]]
    col_name = str(op.get("column", ""))
    if col_name not in headers:
        return f"skipped: column not found ({col_name})"
    col = headers.index(col_name) + 1
    for r in range(2, ws.max_row + 1):
        if str(ws.cell(row=r, column=1).value or "") == str(op.get("req_id", "")):
            ws.cell(row=r, column=col, value=op.get("value"))
            return "applied"
    return f"skipped: req_id row not found ({op.get('req_id')})"


def xlsx_update_cell(wb, p, op):
    sheet = op.get("sheet") or "RTM"
    if sheet not in wb.sheetnames:
        return f"skipped: sheet not found ({sheet})"
    try:
        wb[sheet][str(op.get("cell"))] = op.get("value")
        return "applied"
    except Exception as e:
        return f"skipped: bad cell ({e})"


def pptx_append_changelog_slide(prs, p, op):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = str(op.get("title") or f"변경사항 — {p.get('req_id', '-')}")
    body = None
    for ph in slide.placeholders:
        if ph != slide.shapes.title and ph.has_text_frame:
            body = ph.text_frame
            break
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(9), Inches(5)).text_frame
    bullets = [str(b) for b in (op.get("bullets") or ["(내용 없음)"])]
    body.text = bullets[0]
    for ln in bullets[1:]:
        para = body.add_paragraph()
        para.text = ln
        para.font.size = PPt(14)
    return "applied"


def pptx_replace_text(prs, p, op):
    match = str(op.get("match", ""))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                if "".join(r.text for r in para.runs) == match or para.text == match:
                    for r in list(para.runs):
                        r.text = ""
                    (para.runs[0] if para.runs else para.add_run()).text = str(op.get("new_text", ""))
                    return "applied"
    return "skipped: text not found"


OPS = {
    ("docx", "append_changelog"): ("doc", docx_append_changelog),
    ("docx", "append_after_heading"): ("doc", docx_append_after_heading),
    ("docx", "replace_paragraph"): ("doc", docx_replace_paragraph),
    ("xlsx", "append_row"): ("wb", xlsx_append_row),
    ("xlsx", "update_by_req_id"): ("wb", xlsx_update_by_req_id),
    ("xlsx", "update_cell"): ("wb", xlsx_update_cell),
    ("pptx", "append_changelog_slide"): ("prs", pptx_append_changelog_slide),
    ("pptx", "replace_text"): ("prs", pptx_replace_text),
}


def v1_to_ops(p):
    ads = p.get("affected_deliverables") or [{}]
    return [
        {"target": "docx", "op": "append_changelog", "bullets": p.get("changes", []),
         "risks": p.get("risks", []), "open_questions": p.get("open_questions", []),
         "wiki_commit_url": p.get("wiki_commit_url", "")},
        {"target": "xlsx", "op": "append_row", "sheet": "RTM", "values": [
            p.get("req_id", "-"), p.get("ts", "-"), p.get("requestor", "-"), p.get("source", "-"),
            p.get("decision", "-"), (p.get("summary", "") or "")[:200],
            "; ".join(str(a.get("name", "-")) for a in ads),
            "; ".join(str(a.get("change_type", "-")) for a in ads),
            p.get("da_score", ""), p.get("wiki_commit_url", "")]},
        {"target": "pptx", "op": "append_changelog_slide",
         "title": f"변경사항 — {p.get('req_id', '-')}",
         "bullets": [f"요청자: {p.get('requestor', '-')} · 결정: {p.get('decision', '-')}"]
                    + [f"• {c}" for c in (p.get("changes") or [])[:6]]},
    ]


# ── 상태 추출본 ───────────────────────────────────────────────────────
def export_state(dirpath: Path, doc, wb, prs):
    st = dirpath / "state"
    st.mkdir(parents=True, exist_ok=True)
    lines = []
    for para in doc.paragraphs:
        s, tx = para.style.name, para.text
        if not tx.strip():
            continue
        if s == "Title":
            lines.append(f"# {tx}")
        elif s.lower().startswith("heading"):
            try:
                lv = int(s.split()[-1])
            except Exception:
                lv = 1
            lines.append("#" * min(6, lv + 1) + f" {tx}")
        elif s == "List Bullet":
            lines.append(f"- {tx}")
        else:
            lines.append(tx)
    (st / "docx.md").write_text("\n".join(lines) + "\n", encoding="utf-8")

    buf = io.StringIO()
    w = csv.writer(buf)
    if "RTM" in wb.sheetnames:
        for row in wb["RTM"].iter_rows(values_only=True):
            w.writerow(["" if v is None else v for v in row])
    else:
        w.writerow(RTM_HEADER)
    (st / "rtm.csv").write_text(buf.getvalue(), encoding="utf-8")

    ov = []
    for name in wb.sheetnames:
        ws = wb[name]
        ov.append(f"## Sheet: {name} ({ws.max_row}행 × {ws.max_column}열)")
        for r, row in enumerate(ws.iter_rows(values_only=True), start=1):
            if r > 5:
                ov.append("…")
                break
            ov.append("| " + " | ".join(str(v)[:24] if v is not None else "" for v in row[:8]))
    (st / "xlsx_overview.md").write_text("\n".join(ov) + "\n", encoding="utf-8")

    out = []
    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title is not None else "(제목 없음)"
        out.append(f"## Slide {i}: {title}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        out.append(f"- {para.text}")
    (st / "pptx.md").write_text("\n".join(out) + "\n", encoding="utf-8")


def load_bundle(dirpath: Path):
    files, _ = resolve_files(dirpath)
    return files, {"doc": open_docx(files["docx"]), "wb": open_xlsx(files["xlsx"]), "prs": open_pptx(files["pptx"])}


def save_bundle(files, ctx):
    ctx["doc"].save(str(files["docx"]))
    ctx["wb"].save(str(files["xlsx"]))
    ctx["prs"].save(str(files["pptx"]))


def export_state_all():
    n = 0
    cand = []
    if DELIV.exists():
        cand += [p for p in DELIV.iterdir() if p.is_dir() and p.name not in ("patches",)]
    scen = ROOT / "scenarios"
    if scen.exists():
        cand += [p / "deliverables" for p in scen.iterdir() if (p / "deliverables").is_dir()]
    for d in sorted(cand):
        if not ((d / "manifest.json").exists() or any(discover(d, e) for e in ("docx", "xlsx", "pptx"))):
            continue
        files, ctx = load_bundle(d)
        export_state(d, ctx["doc"], ctx["wb"], ctx["prs"])
        n += 1
        print(f"state exported: {d.relative_to(ROOT)}")
    print(f"export-state-all: {n} dir(s)")


def main() -> int:
    DELIV.mkdir(parents=True, exist_ok=True)
    if "--export-state-all" in sys.argv:
        export_state_all()
        return 0

    PENDING.mkdir(parents=True, exist_ok=True)
    APPLIED.mkdir(parents=True, exist_ok=True)
    patches = sorted(PENDING.glob("patch_*.json"))
    if not patches:
        print("no pending patches")
        return 0

    bundles = {}   # dir → (files, ctx)  — 여러 패치가 같은 디렉토리를 공유해도 1회 로드
    done = []
    for path in patches:
        try:
            p = json.loads(path.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"skip (invalid json): {path.name}: {e}")
            continue
        schema = p.get("schema")
        if schema == "deliverable_patch/v2":
            ops = (p.get("ops") or [])[:MAX_OPS]
            dirpath = ROOT / str(p.get("deliverable_dir") or "deliverables/live")
        elif schema == "deliverable_patch/v1":
            ops = v1_to_ops(p)
            dirpath = DELIV                                   # v1 루트 파일 하위호환
        else:
            print(f"skip (unknown schema): {path.name}")
            continue
        _ok = any(str(dirpath.resolve()).startswith(str((ROOT / pfx).resolve()))
                  for pfx in ("scenarios", "deliverables")) or dirpath == DELIV
        if not _ok:
            print(f"skip (dir outside scenarios//deliverables/): {path.name}")
            continue
        if dirpath not in bundles:
            bundles[dirpath] = load_bundle(dirpath)
        files, ctx = bundles[dirpath]

        results = []
        for op in ops:
            key = (op.get("target"), op.get("op"))
            if key not in OPS:
                results.append({"op": op, "result": f"skipped: unknown op {key}"})
                continue
            kind, fn = OPS[key]
            try:
                results.append({"op": op, "result": fn(ctx[kind], p, op)})
            except Exception as e:
                results.append({"op": op, "result": f"skipped: error {e}"})
        done.append((path, results))
        print(f"processed {path.name} → {dirpath.name}: "
              f"{sum(1 for r in results if r['result'] == 'applied')}/{len(results)} ops applied")

    for dirpath, (files, ctx) in bundles.items():
        save_bundle(files, ctx)
        export_state(dirpath if dirpath != DELIV else DELIV, ctx["doc"], ctx["wb"], ctx["prs"])
    for path, results in done:
        (APPLIED / (path.stem + ".result.json")).write_text(
            json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
        shutil.move(str(path), str(APPLIED / path.name))
    print(f"done: {len(done)} patch(es)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
